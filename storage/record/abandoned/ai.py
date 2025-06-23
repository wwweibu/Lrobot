# ai 相关逻辑
# 使用ai时需要更改数据库展示页面，处理blob类型数据，需要修改 get_database函数中：all_data[table] = [serialize_row(row) for row in rows]
def serialize_row(row: dict) -> dict:
    """序列化数据库行，处理 BLOB 字段为 base64 字符串"""
    new_row = {}
    for key, value in row.items():
        if isinstance(value, bytes):  # 处理 BLOB 类型
            new_row[key] = base64.b64encode(value).decode("ascii")
        else:
            new_row[key] = value
    return new_row

# aichat.py ai对话功能，放在command下
#实现ai对话逻辑，由aichat.py调用chat.py对话循环，chat.py调用model.py与模型通信，同时调用tool.py生成函数的schema传入模型
from message.handler.msg import Msg
from ..infra import add_status, delete_status, update_database, identify_user


async def aichat_open(msg: Msg):
    qq = await identify_user(msg.source, msg.robot)
    if not qq:
        content = "开启失败"
    else:
        content = "开启成功"
        await add_status(qq, "对话")
    kind = msg.kind[:2]
    Msg(
        robot=msg.robot,
        kind=f"{kind}发送文本",
        event="发送",
        source=msg.source,
        seq=msg.seq,
        content=content,
        group=msg.group,
    )


async def aichat_close(msg: Msg):
    kind = msg.kind[:2]
    qq = await identify_user(msg.source, msg.robot)
    await delete_status(qq, "对话")
    await add_status(qq, "对话保存")
    Msg(
        robot=msg.robot,
        kind=f"{kind}发送文本",
        event="发送",
        source=msg.source,
        seq=msg.seq,
        content="请输入/保存来保存对话，/不保存来删除对话",
        group=msg.group,
    )


async def aichat_save_yes(msg: Msg):
    kind = msg.kind[:2]
    qq = await identify_user(msg.source, msg.robot)
    await delete_status(qq, "对话保存")
    Msg(
        robot=msg.robot,
        kind=f"{kind}发送文本",
        event="发送",
        source=msg.source,
        seq=msg.seq,
        content="对话已保存",
        group=msg.group,
    )


async def aichat_save_no(msg: Msg):
    kind = msg.kind[:2]
    qq = await identify_user(msg.source, msg.robot)
    await delete_status(qq, "对话保存")
    query = "DELETE FROM logic_chat WHERE role = ?"
    params = (qq,)
    await update_database(query, params)
    Msg(
        robot=msg.robot,
        kind=f"{kind}发送文本",
        event="发送",
        source=msg.source,
        seq=msg.seq,
        content="对话已删除",
        group=msg.group,
    )

# 对话（chat文件夹）
# ai对话循环chat.py
import json
from .model import model_request
from log import loggers
from message.handler.msg import Msg
from ..infra import query_database, update_database, identify_user
from .tool import get_tools

index = 0
model_list = ["ollama_qwen", "silicon", "zhipu", "hunyuan", "huoshan", "spark"]


async def chat(msg: Msg):
    """ai对话"""
    global index
    kind = msg.kind[:2]
    query = "SELECT model, message FROM logic_chat WHERE role = ?"
    qq = await identify_user(msg.source, msg.robot)
    result = await query_database(query, (qq,))

    if not result:
        model = model_list[index % len(model_list)]

        # 插入新角色记录
        await update_database(
            "INSERT INTO logic_chat (role, model) VALUES (?, ?)",
            (qq, model),
        )
    tool_list, func_map = get_tools(cab=0)
    context = "你是一个智能助手，当你需要从工具中选择并调用必要工具时，要根据问题需求进行精准筛选。工具可以分多次、多轮调用，直到你获得想要的结果。特别地，你每轮对话中必须调用知识库工具search_kb，需从问题中提取关键信息并生成精准的搜索词列表，再通过调用知识库工具进行搜索，参数为[a,b,c]，以获取相关知识数据。"
    context1 = "你是LRobot，武汉大学逻辑推理协会的社团小助手"
    response = await chat_loop(qq, msg.content, tool_list, func_map, context, context1)
    Msg(
        robot=msg.robot,
        kind=f"{kind}发送文本",
        event="发送",
        source=msg.source,
        seq=msg.seq,
        content=response,
        group=msg.group,
    )


async def chat_loop(
    role, content, tool_list=None, tool_map=None, prompt1=None, prompt2=None
):
    """对话循环"""
    query = "SELECT model,message FROM logic_chat WHERE role = ?"
    result = await query_database(query, (role,))
    if result:
        model = result[0]["model"]
        message = json.loads(result[0]["message"] or "[]")
    else:
        model = "hunyuan"
        message = []
    if prompt1:
        message = insert_system_message(message, prompt1)
    message.append({"role": "user", "content": content})

    response_message = await model_request(model, message, tool_list)
    print(response_message["content"])
    message.append(response_message)

    while response_message.get("tool_calls"):
        for tool_call in response_message["tool_calls"]:
            if function_to_call := tool_map.get(tool_call["function"]["name"]):
                arguments = tool_call["function"]["arguments"]
                if isinstance(arguments, str):
                    arguments = json.loads(arguments)
                try:
                    output = await function_to_call(**arguments)
                    loggers["message"].info(
                        f"函数调用 -> 函数: {tool_call['function']['name']} | 参数: {arguments}",
                        extra={"event": "对话消息"},
                    )
                except Exception as e:
                    output = e
            else:
                loggers["message"].error(
                    f"函数调用异常 -> 不存在函数{tool_call['function']['name']}",
                    extra={"event": "对话消息"},
                )
                output = "function not found"
            if "id" in tool_call:
                result = {
                    "role": "tool",
                    "content": str(output),
                    "tool_call_id": tool_call["id"],
                }
            else:
                result = {
                    "role": "tool",
                    "content": str(output),
                    "name": tool_call["function"]["name"],
                }
            message.append(result)
            print(result)
            response_message = await model_request(model, message, tool_list)
            message.append(response_message)
            print(response_message["content"])

    result_message = message.copy()
    if prompt2:
        result_message = insert_system_message(message, prompt2)

    result_message.append(
        {
            "role": "user",
            "content": f"现在，你需要结合上下文来重新回答这个问题：{content}。如果问题能够得到解决，你要提出合理的建议并回复；如果未解决，则回复'否'",
        }
    )
    response_message = await model_request(model, result_message)
    print(response_message["content"])
    print(result_message)
    result_message.append(response_message)
    query = """
            INSERT INTO logic_chat (role,message) VALUES (?,?)
            ON CONFLICT(role) DO UPDATE SET message = ?
            """
    params = (role, json.dumps(message), json.dumps(message))
    await update_database(query, params)
    return response_message["content"]


def insert_system_message(message, system_content):
    """插入 system 提示"""
    system_message = {"role": "system", "content": system_content}
    if message:
        if message[0].get("role") == "system":
            message[0] = system_message
        else:
            message.insert(0, system_message)
    else:
        message = [system_message]

    return message

# 知识库生成kb.py
import os
import xlrd
import chardet
import whisper
import rarfile
import zipfile
import aiofiles
import openpyxl
import textract
import pythoncom
import pytesseract
import numpy as np
from PIL import Image
import win32com.client
from docx import Document
import moviepy.editor as mp3
from numpy.linalg import norm
from bs4 import BeautifulSoup
from pptx import Presentation
from moviepy.editor import AudioFileClip
from sentence_transformers import SentenceTransformer
from config import path, clear_socks
from log import loggers
from ..infra import query_database, update_database

system_logger = loggers["system"]
prompt = "为这个句子生成表示以用于检索："

SUPPORTED_EXTS = {
    ".txt",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".ppt",
    ".jpg",
    ".png",
    ".zip",
    ".rar",
    ".mp4",
    ".wmv",
    ".html",
    ".mp3",
    ".doc",
}


async def extract_text_from_doc(file_path):
    """从DOC文件中提取文本"""
    pythoncom.CoInitialize()
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False  # 不显示 Word 窗口

    # 打开文件
    doc = word.Documents.Open(file_path)

    # 提取文本
    text = doc.Content.Text

    # 关闭文档并退出 Word
    doc.Close()
    word.Quit()

    return text.strip()


async def extract_text_from_txt(file_path):
    """从TXT文件中提取文本"""
    with open(file_path, "rb") as f:
        raw_data = f.read(10000)
        detected = chardet.detect(raw_data)
        encoding = detected["encoding"] or "utf-8"

    # 使用异步方式读取全文
    async with aiofiles.open(
        file_path, mode="r", encoding=encoding, errors="ignore"
    ) as f:
        text = await f.read()

    return text.strip()


async def extract_text_from_html(file_path):
    """从HTML文件中提取文本"""
    with open(file_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")
        text = soup.get_text()
    return text.strip()


async def extract_text_from_xls(file_path):
    """从XLS文件中提取文本"""
    book = xlrd.open_workbook(file_path)
    sheet = book.sheet_by_index(0)
    text = ""
    for row in range(sheet.nrows):
        text += (
            " ".join(str(sheet.cell_value(row, col)) for col in range(sheet.ncols))
            + "\n"
        )
    return text.strip()


async def extract_text_from_mp3(file_path):
    """从MP3文件中提取文本（使用 whisper）"""
    wav_path = "temp_audio.wav"
    audio_clip = AudioFileClip(file_path)
    audio_clip.write_audiofile(wav_path, codec="pcm_s16le")

    whisper_model = whisper.load_model("base")
    result = whisper_model.transcribe(wav_path, language="zh")

    os.remove(wav_path)
    return result["text"]


async def extract_text_from_video(file_path):
    """从视频文件中提取文本（语音转文本）"""
    audio = mp3.VideoFileClip(file_path).audio
    audio.write_audiofile("audio.wav")
    whisper_model = whisper.load_model("base")
    result = whisper_model.transcribe("audio.wav", language="zh")
    os.remove("audio.wav")
    return result["text"]


async def extract_text_from_image(file_path):
    """从图片文件中提取文本"""
    img = Image.open(file_path)
    img = pytesseract.image_to_string(img, lang="chi_sim")  # 提取中文文本
    return img.strip()


async def extract_text_from_pptx(file_path):
    """从PPT文件中提取文本"""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()


async def extract_text_from_docx(file_path):
    """从Word文件中提取文本"""
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text.strip()


async def extract_text_from_xlsx(file_path):
    """从Excel文件中提取文本"""
    wb = openpyxl.load_workbook(file_path)
    text = ""
    for sheet in wb.sheetnames:
        sheet_data = wb[sheet]
        for row in sheet_data.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row]) + "\n"
    return text.strip()


async def extract_text_from_zip(file_path):
    """从ZIP文件中提取文本"""
    with zipfile.ZipFile(file_path, "r") as zip_ref:
        text = ""
        for file_name in zip_ref.namelist():
            if file_name.endswith(".txt"):
                with zip_ref.open(file_name) as file:
                    text += file.read().decode("utf-8") + "\n"
        return text.strip()


async def extract_text_from_rar(file_path):
    """从RAR文件中提取文本"""
    with rarfile.RarFile(file_path) as rar_ref:
        text = ""
        for file_name in rar_ref.namelist():
            if file_name.endswith(".txt"):
                with rar_ref.open(file_name) as file:
                    text += file.read().decode("utf-8") + "\n"
        return text.strip()


async def extract_text_from_pdf(file_path):
    """从PDF文件中提取文本"""
    text = textract.process(file_path).decode("utf-8", errors="ignore")
    return text.strip()


@clear_socks
async def process_kb():
    """处理文件夹并更新知识库数据库"""
    model = SentenceTransformer("BAAI/bge-large-zh-v1.5")
    folder = path / "storage/file/resource/kb"
    files_to_update = []
    # 获取数据库中已有的记录
    existing = await query_database("SELECT path, mtime FROM file_kb")
    db_index = {row["path"]: row["mtime"] for row in existing}

    for file in folder.rglob("*"):
        if not file.is_file():  # 跳过文件夹
            continue
        file_path = str(file.resolve())
        mtime = os.path.getmtime(file_path)
        text = ""

        # 跳过未更新文件
        if file_path in db_index and mtime <= db_index[file_path]:
            continue
        if file.suffix.lower() not in SUPPORTED_EXTS:
            system_logger.debug(f"不支持的文件 {file}", extra={"event": "知识更新"})
            files_to_update.append(
                {
                    "file_path": file_path,
                    "mtime": mtime,
                    "text": "",
                    "vector": None,
                }
            )
            continue

        system_logger.debug(f"文件{file}正在更新……", extra={"event": "知识更新"})
        try:
            if file.suffix.lower() == ".jpg" or file.suffix.lower() == ".png":
                text = await extract_text_from_image(file_path)
            elif file.suffix.lower() == ".txt":
                text = await extract_text_from_txt(file_path)
            elif file.suffix.lower() == ".html":
                text = await extract_text_from_html(file_path)
            elif file.suffix.lower() == ".docx":
                text = await extract_text_from_docx(file_path)
            elif file.suffix.lower() == ".doc":
                text = await extract_text_from_doc(file_path)
            elif file.suffix.lower() == ".xls":
                text = await extract_text_from_xls(file_path)
            elif file.suffix.lower() == ".xlsx":
                text = await extract_text_from_xlsx(file_path)
            elif file.suffix.lower() == ".pptx":
                text = await extract_text_from_pptx(file_path)
            elif file.suffix.lower() == ".ppt":
                text = await extract_text_from_pptx(file_path)
            elif file.suffix.lower() == ".pdf":
                text = await extract_text_from_pdf(file_path)
            elif file.suffix.lower() == ".zip":
                text = await extract_text_from_zip(file_path)
            elif file.suffix.lower() == ".rar":
                text = await extract_text_from_rar(file_path)
            elif file.suffix.lower() == ".mp4":
                text = await extract_text_from_video(file_path)
            elif file.suffix.lower() == ".wmv":
                text = await extract_text_from_video(file_path)
            elif file.suffix.lower() == ".mp3":
                text = await extract_text_from_mp3(file_path)

            if not text:
                system_logger.debug(f"不支持的文件 {file}", extra={"event": "知识更新"})
                files_to_update.append(
                    {
                        "file_path": file_path,
                        "mtime": mtime,
                        "text": "",
                        "vector": None,
                    }
                )
                continue

            text = f"文件名：{file_path.replace(str(folder), '')}\n内容：{text}"
            title_vec = model.encode(
                ["为这个标题生成表示以用于检索：" + file_path.replace(str(folder), "")],
                normalize_embeddings=True,
            )[0]
            content_vec = model.encode(
                ["为这个正文生成表示以用于检索：" + text], normalize_embeddings=True
            )[0]
            final_vec = 0.3 * title_vec + 0.7 * content_vec
            files_to_update.append(
                {
                    "file_path": file_path,
                    "mtime": mtime,
                    "text": text,
                    "vector": np.array(final_vec, dtype=np.float32),
                }
            )
        except Exception as e:
            system_logger.error(
                f"跳过文件: {file_path} | 原因: {e}", extra={"event": "知识更新"}
            )
            files_to_update.append(
                {
                    "file_path": file_path,
                    "mtime": mtime,
                    "text": "",
                    "vector": None,
                }
            )

    print(files_to_update)
    print("正在写入数据库……")

    if files_to_update:
        query = """
                                INSERT INTO file_kb (path, mtime, content, vector)
                                VALUES (?, ?, ?, ?)
                                ON CONFLICT(path) DO UPDATE SET 
                                    mtime = excluded.mtime,
                                    content = excluded.content,
                                    vector = excluded.vector
                                """
        for file in files_to_update:
            params = (
                file["file_path"],
                file["mtime"],
                file["text"],
                (
                    file["vector"].tobytes() if file["vector"] is not None else None
                ),  # FAISS向量存储为 BLOB
            )
            await update_database(query, params)
            system_logger.debug(
                f"文件更新成功: {file['file_path']}", extra={"event": "知识更新"}
            )


def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    return float(np.dot(a, b) / (norm(a) * norm(b)))


@clear_socks
async def search_kb(queries: list):
    """
    在社团知识库里搜索，需要简化问题生成对应的关键词，返回匹配结果列表
    @param queries: 关键词列表
    """
    model = SentenceTransformer("BAAI/bge-large-zh-v1.5")
    query = "SELECT path, mtime, content, vector FROM file_kb"
    results = await query_database(query, ())

    kb_entries = []
    for result in results:
        if result["content"]:
            vector = np.frombuffer(result["vector"], dtype=np.float32)
            kb_entries.append(
                {
                    "path": result["path"],
                    "mtime": result["mtime"],
                    "content": result["content"],
                    "vector": vector,
                }
            )

    # 查询处理
    results = []
    for query in queries:
        q_vec = model.encode(
            [f"为这个句子生成表示以用于检索：{query}"], normalize_embeddings=True
        )[0]
        matched = []
        for entry in kb_entries:
            score = cosine_similarity(q_vec, entry["vector"])
            print(score)
            if score >= 0.5:
                matched.append(
                    {
                        "query": query,
                        "content": entry["content"],
                        "path": entry["path"],
                        "score": round(score, 3),
                    }
                )
        results.append(matched)

    return results

# 模型调用model.py
# 需要在model.yaml中输入以下数据：Zhipu_KEY:Spark_URL:Spark_KEY:Hunyuan_URL:Hunyuan_KEY:Huoshan_API:Huoshan_KEY:Huoshan_URL:Silicon_URL:Silicon_KEY:
import json
import httpx
import asyncio
from openai import OpenAI
from zhipuai import ZhipuAI
from starlette.websockets import WebSocketDisconnect
from config import clear_socks, config

connected_service = {}


def set_connected_service(name, websocket):
    connected_service[name] = websocket


@clear_socks
async def model_request(model: str, messages: list, tool=None):
    """根据模型调用不同的 API 获取响应"""
    if model.startswith("ollama"):
        if connected_service.get("ollama") is None:
            raise Exception("模型服务未连接")
        request = {
            "model": model,
            "messages": messages,
            "tool": tool,
        }
        await connected_service["ollama"].send_json(request)
        try:
            while True:
                response = await asyncio.wait_for(
                    connected_service["ollama"].receive_json(), timeout=120
                )
                if response.get("type") == "heartbeat":
                    print("收到心跳包，继续等待...")
                    continue  # 继续等待实际响应
                return response  # 实际响应数据
        except WebSocketDisconnect as e:
            print(f"连接异常断开: Code={e.code}, Reason={e.reason}")
            raise
        except asyncio.TimeoutError:
            raise Exception("模型响应超时，请检查服务是否正常运行")
    elif model == "hunyuan":
        client = OpenAI(
            api_key=config["Hunyuan_KEY"],
            base_url=config["Hunyuan_URL"],
        )
        completion = client.chat.completions.create(
            model="hunyuan-functioncall",
            messages=messages,
            tools=tool,
        )
        response = completion.model_dump()
        return response["choices"][0]["message"]
    elif model == "huoshan":
        client = OpenAI(
            api_key=config["Huoshan_API"],
            base_url=config["Huoshan_URL"],
        )
        completion = client.chat.completions.create(
            model="doubao-pro-32k-functioncall-241028",
            messages=messages,
            tools=tool,
        )
        response = completion.model_dump()
        return response["choices"][0]["message"]
    elif model == "zhipu":
        client = ZhipuAI(api_key=config["Zhipu_KEY"])
        completion = client.chat.completions.create(
            model="glm-4-flash",
            messages=messages,
            tools=tool,
        )
        response = completion.model_dump()
        return response["choices"][0]["message"]
    async with httpx.AsyncClient(timeout=httpx.Timeout(15.0)) as client:
        if model == "spark":
            url = config["Spark_URL"]
            headers = {"Authorization": f"Bearer {config['Spark_KEY']}"}
            data = {
                "model": "4.0Ultra",
                "messages": messages,
                **({"tools": tool} if tool is not None else {}),
            }
            response = await client.post(url, headers=headers, json=data)
            response = json.loads(response.text)
            return response["choices"][0]["message"]
        elif model == "silicon":
            url = config["Silicon_URL"]
            headers = {
                "Authorization": f"Bearer {config['Silicon_KEY']}",
                "Content-Type": "application/json",
            }
            data = {
                "model": "THUDM/GLM-4-9B-0414",
                "messages": messages,
                **({"tools": tool} if tool is not None else {}),
            }
            response = await client.post(url, headers=headers, json=data)
            response = json.loads(response.text)
            return response["choices"][0]["message"]

#函数调用tool.py
# 函数调用
import re
import time
import inspect
from typing import Callable
from .kb import search_kb


def parse_docstring(docstring: str):
    """
    解析函数的 docstring，提取 @param、@enum、@description 信息
    """
    param_pattern = re.compile(r"@param (\w+): (.+)")
    enum_pattern = re.compile(r"@enum: \[([^\]]+)\]")
    description_pattern = re.compile(r"@description: (.+)")
    format_pattern = re.compile(r"@format: (.+)")
    min_length_pattern = re.compile(r"@minLength: (\d+)")
    max_length_pattern = re.compile(r"@maxLength: (\d+)")

    params_info = {}

    lines = docstring.split("\n")
    current_param = None

    for line in lines:
        line = line.strip()

        # 解析 @param
        param_match = param_pattern.match(line)
        if param_match:
            current_param, desc = param_match.groups()
            params_info[current_param] = {"description": desc}
            continue

        # 解析 @enum
        enum_match = enum_pattern.search(line)
        if enum_match and current_param:
            params_info[current_param]["enum"] = [
                eval(x.strip()) for x in enum_match.group(1).split(",")
            ]

        # 解析 @description
        desc_match = description_pattern.match(line)
        if desc_match and current_param:
            params_info[current_param]["description"] = desc_match.group(1)

        # 解析 @format
        format_match = format_pattern.match(line)
        if format_match and current_param:
            params_info[current_param]["format"] = format_match.group(1)

        # 解析 @minLength
        min_length_match = min_length_pattern.match(line)
        if min_length_match and current_param:
            params_info[current_param]["minLength"] = int(min_length_match.group(1))

        # 解析 @maxLength
        max_length_match = max_length_pattern.match(line)
        if max_length_match and current_param:
            params_info[current_param]["maxLength"] = int(max_length_match.group(1))

    return params_info


def function_to_json_schema(func):
    """
    将 Python 函数解析为 JSON Schema 格式。
    """
    func_name = func.__name__
    docstring = inspect.getdoc(func) or ""  # 文档字符串
    params_info = parse_docstring(docstring)  # 参数注释

    # 获取函数参数
    sig = inspect.signature(func)
    required_params = []
    properties = {}

    for param_name, param in sig.parameters.items():
        param_schema = {}

        # 提取类型信息
        if param.annotation is int:
            param_schema["type"] = "integer"
        elif param.annotation is str:
            param_schema["type"] = "string"
        elif param.annotation is float:
            param_schema["type"] = "number"
        elif param.annotation is bool:
            param_schema["type"] = "boolean"
        elif param.annotation is list:
            param_schema["type"] = "array"
        elif param.annotation is dict:
            param_schema["type"] = "object"
        else:
            param_schema["type"] = "string"  # 默认字符串类型

        # 处理默认值
        if param.default is not param.empty:
            param_schema["default"] = param.default
        else:
            required_params.append(param_name)

        # 处理注释信息
        if param_name in params_info:
            param_schema.update(params_info[param_name])

        properties[param_name] = param_schema

    # 生成最终 JSON Schema
    schema = {
        "type": "function",
        "function": {
            "name": func_name,
            "description": docstring.split("\n")[0] if docstring else "",
            "parameters": {
                "type": "object",
                "required": required_params,
                "properties": properties,
            },
        },
    }

    return schema


# 函数示例
async def get_current_time():
    """
    用于获取当前时间
    """
    return int(time.time())


# 函数示例
async def get_symbol(txt: str):
    """
    使用尼古拉加密法加密字符串
    @param txt: 需要加密的字符串
        @maxLength: 3
        @enum: ["ddd","efg","abc"]
    """
    return "xyz"


def build_tools(*funcs: Callable):
    """生成工具列表"""
    tool_list = []
    func_map = {}
    for func in funcs:
        if not inspect.iscoroutinefunction(func):  # 必须为异步函数
            raise TypeError(
                f"函数 {func.__name__} 不是异步函数（async def），请修改为异步形式"
            )
        schema = function_to_json_schema(func)
        tool_list.append(schema)
        func_map[schema["function"]["name"]] = func
    return tool_list, func_map


def get_tools(cab: int = None, func_list: list = None):
    """获取工具列表"""
    # TODO 修改工具列表
    if cab:
        func_list = [get_current_time, get_symbol]
    else:
        func_list = [get_symbol, get_current_time, search_kb]
    if func_list:
        func_list = func_list
    tool_list, func_map = build_tools(*func_list)
    return tool_list, func_map




# 连接
# model_connect.py 服务端接口
import asyncio
from fastapi import WebSocket, APIRouter
from logic import set_connected_service

router = APIRouter()


@router.websocket("/ollama")
async def websocket_endpoint(websocket: WebSocket):
    await websocket.accept()
    set_connected_service("ollama", websocket)
    print("服务方已连接")

    try:
        while True:
            await asyncio.sleep(1)  # 保持连接，不主动关闭
    except Exception as e:
        print("连接断开:", e)
        set_connected_service("ollama", None)

# model_join.py 客户端接口
import json
import asyncio
import websockets
from ollama import chat, ChatResponse


def ollama_model_request(model, messages, tool):
    # 模拟一个长时间任务
    model = model.replace("ollama_", "")
    if model == "qwen":
        model = "qwen2.5:14b"
    completion: ChatResponse = chat(model, messages=messages, tools=tool)
    response = completion.model_dump()
    return response["message"]


async def send_heartbeat(websocket, interval=5):
    """每隔 interval 秒发送一次心跳状态"""
    try:
        while True:
            await asyncio.sleep(interval)
            await websocket.send(json.dumps({"status": "working", "type": "heartbeat"}))
    except websockets.exceptions.ConnectionClosed:
        # 连接已关闭，终止心跳
        pass


async def serve_model(uri: str):
    async with websockets.connect(uri) as websocket:
        print("连接到调用方成功")

        while True:
            try:
                message = await websocket.recv()
                request = json.loads(message)
                print("收到请求:", request)

                model = request.get("model", "")
                messages = request.get("messages", [])
                tool = request.get("tool", None)

                # 创建心跳发送任务
                task_heartbeat = asyncio.create_task(send_heartbeat(websocket))

                # 创建模型处理任务
                result = await asyncio.to_thread(
                    ollama_model_request, model, messages, tool
                )

                # 等待模型处理完成
                print("模型处理完成:", result)

                # 模型完成后取消心跳任务
                task_heartbeat.cancel()
                try:
                    await task_heartbeat
                except asyncio.CancelledError:
                    pass

                # 发送最终结果
                await websocket.send(json.dumps(result))

            except Exception as e:
                print("处理失败，断开连接:", e)
                break


# 启动客户端
if __name__ == "__main__":
    asyncio.run(serve_model("wss://whumystery.cn/hjd/ollama"))
